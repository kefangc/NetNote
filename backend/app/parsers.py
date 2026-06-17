from __future__ import annotations

import re
from pathlib import Path


def parse_file(path: Path, filename: str) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".markdown", ".csv"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp"}:
        return _parse_image_placeholder(filename)
    return path.read_text(encoding="utf-8", errors="ignore")


def _parse_pdf(path: Path) -> str:
    try:
        import pdfplumber

        pages: list[str] = []
        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[第 {index} 页]\n{text}")
        return "\n\n".join(pages) or "PDF 未提取到可读文本，可在 V2 接入 OCR 管线补充扫描件识别。"
    except Exception as exc:
        return f"PDF 解析失败：{exc}"


def _parse_docx(path: Path) -> str:
    try:
        from docx import Document

        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as exc:
        return f"DOCX 解析失败：{exc}"


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        slides: list[str] = []
        for index, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[第 {index} 页幻灯片]\n" + "\n".join(texts))
        return "\n\n".join(slides) or "PPTX 未提取到可读文本。"
    except Exception as exc:
        return f"PPTX 解析失败：{exc}"


def _parse_image_placeholder(filename: str) -> str:
    return (
        f"图片来源 {filename} 已登记。当前 V1 保留 OCR 管线入口；"
        "如安装 PaddleOCR，可在 parsers.py 中替换该占位解析为真实 OCR 文本。"
    )


def chunk_text(text: str, source_id: str, source_title: str, size: int = 720, overlap: int = 120):
    from .schemas import SourceChunk
    from .agents import extract_keywords, make_id

    clean = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not clean:
        clean = "该来源暂无可读文本。"
    chunks: list[SourceChunk] = []
    start = 0
    index = 1
    while start < len(clean):
        end = min(len(clean), start + size)
        segment = clean[start:end].strip()
        if segment:
            chunks.append(
                SourceChunk(
                    id=make_id("chunk"),
                    source_id=source_id,
                    source_title=source_title,
                    text=segment,
                    location=f"片段 {index}",
                    keywords=extract_keywords(segment),
                )
            )
        if end == len(clean):
            break
        start = max(end - overlap, start + 1)
        index += 1
    return chunks
